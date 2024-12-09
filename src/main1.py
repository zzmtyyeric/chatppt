from pptx.dml.color import RGBColor
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE  # 导入图表类型
import numpy as np

def haha():
    presentation = Presentation("outputs/ChatPPT Demo.pptx")
    # 添加文本内容幻灯片
    slide_layout = presentation.slide_layouts[-1]
    slide = presentation.slides.add_slide(slide_layout)


    # 打印新增页属性
    print(f"Slide ID: {slide.slide_id}")
    print(f"  Layout: {slide.slide_layout}")
    print(f"  Shapes: {len(slide.shapes)} shapes")
    print(f"  Placeholders: {len(slide.placeholders)} placeholders")

    print("  Shape Details:")
    for shape in slide.shapes:
        print(f"    - Shape Name: {shape.name}, Type: {shape.shape_type}")

    print("  Placeholder Details:")
    for placeholder in slide.placeholders:
        print(f"    - Placeholder ID: {placeholder.placeholder_format.idx}, Type: {placeholder.placeholder_format.type}")

    # 填充原有布局中的占位符（标题和文本）
    title = slide.shapes.title
    title.text = "python-pptx 新增文本内容示例"
    content = slide.placeholders[14]
    content.text = "填充原有的文本占位符"

def addText():
    presentation = Presentation("outputs/ChatPPT Demo.pptx")
    # 添加文本内容幻灯片
    slide_layout = presentation.slide_layouts[-1]
    slide = presentation.slides.add_slide(slide_layout)


    # 打印新增页属性
    print(f"Slide ID: {slide.slide_id}")
    print(f"  Layout: {slide.slide_layout}")
    print(f"  Shapes: {len(slide.shapes)} shapes")
    print(f"  Placeholders: {len(slide.placeholders)} placeholders")

    print("  Shape Details:")
    for shape in slide.shapes:
        print(f"    - Shape Name: {shape.name}, Type: {shape.shape_type}")

    print("  Placeholder Details:")
    for placeholder in slide.placeholders:
        print(f"    - Placeholder ID: {placeholder.placeholder_format.idx}, Type: {placeholder.placeholder_format.type}")

    # 填充原有布局中的占位符（标题和文本）
    content = slide.placeholders[10]
    content.text = "python-pptx 新增文本内容示例"

    # 新增文本框
    left = Inches(6)
    top = Inches(5)
    width = Inches(5)
    height = Inches(1)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = "额外的文本框内容"

    # 格式化文本
    paragraph = text_frame.add_paragraph()  # 添加新段落
    paragraph.text = "这是一个新的段落。"  # 设置段落文本

    # 设置字体
    run = paragraph.add_run()  # 添加文本运行
    run.text = " 这部分是加粗的文本。"  # 设置文本内容
    run.font.bold = True  # 设置为粗体
    run.font.size = Pt(16)  # 设置字体大小
    run.font.color.rgb = RGBColor(255, 0, 0)  # 设置字体颜色为红色

    return presentation

def addAll():
    # 创建演示文稿对象
    prs = Presentation()

    # 添加一张幻灯片
    slide_layout = prs.slide_layouts[5]  # 使用空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 1. 添加文本
    text_left = Inches(1)
    text_top = Inches(1)
    text_width = Inches(8)
    text_height = Inches(1)
    text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
    text_frame = text_box.text_frame
    p = text_frame.add_paragraph()
    p.text = "这是一个包含文本、图片、表格和图表的示例"
    p.font.size = Pt(24)

    # 2. 添加图片
    img_path = r'C:\Users\ZZM\Desktop\ZZM幻灯片培训\蜡3.jpg'  # 请替换为你的图片路径
    img_left = Inches(1)
    img_top = text_top + text_height + Inches(0.5)  # 图片顶部距离文本框下方0.5英寸
    slide.shapes.add_picture(img_path, img_left, img_top, width=Inches(4))

    # 3. 添加表格
    table_rows, table_cols = 3, 3  # 表格行数和列数
    table_left = Inches(1)
    table_top = img_top + Inches(2)  # 表格顶部距离图片下方2英寸
    table_width = Inches(6)
    table_height = Inches(2)
    table = slide.shapes.add_table(table_rows, table_cols, table_left, table_top, table_width, table_height).table

    # 填充表格数据
    for i in range(table_rows):
        for j in range(table_cols):
            table.cell(i, j).text = f"单元格 {i+1},{j+1}"

    # 4. 添加图表
    chart_data = CategoryChartData()
    chart_data.categories = ['类别 1', '类别 2', '类别 3']
    chart_data.add_series('系列 1', (19.2, 21.4, 16.7))
    chart_data.add_series('系列 2', (16.5, 22.3, 20.0))

    # 添加图表到幻灯片
    chart_x = Inches(1)
    chart_y = table_top + table_height + Inches(0.5)  # 图表顶部距离表格下方0.5英寸
    chart_cx = Inches(8)  # 图表宽度
    chart_cy = Inches(3)  # 图表高度
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_x, chart_y, chart_cx, chart_cy, chart_data)

    # 保存演示文稿
    prs.save('example_presentation.pptx')

    print("演示文稿已成功创建！")

if __name__ == "__main__":

    p1 = addAll()