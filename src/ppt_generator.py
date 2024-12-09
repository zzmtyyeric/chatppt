import os
from pptx import Presentation
from utils import remove_all_slides
from logger import LOG  # 引入日志模块

# 生成 PowerPoint 演示文稿
def generate_presentation(powerpoint_data, template_path: str, output_path: str):
    # 检查模板文件是否存在
    if not os.path.exists(template_path):
        LOG.error(f"模板文件 '{template_path}' 不存在。")  # 记录错误日志
        raise FileNotFoundError(f"模板文件 '{template_path}' 不存在。")

    prs = Presentation(template_path)  # 加载 PowerPoint 模板
    remove_all_slides(prs)  # 清除模板中的所有幻灯片
    prs.core_properties.title = powerpoint_data.title  # 设置 PowerPoint 的核心标题

    # 遍历所有幻灯片数据，生成对应的 PowerPoint 幻灯片
    for slide in powerpoint_data.slides:
        # 确保布局索引不超出范围，超出则使用默认布局
        if slide.layout_id >= len(prs.slide_layouts):
            slide_layout = prs.slide_layouts[0]
        else:
            slide_layout = prs.slide_layouts[slide.layout_id]

        new_slide = prs.slides.add_slide(slide_layout)  # 添加新的幻灯片

        # 设置幻灯片标题
        if new_slide.shapes.title:
            new_slide.shapes.title.text = slide.content.title
            LOG.debug(f"设置幻灯片标题: {slide.content.title}")

        # 添加文本内容
        for shape in new_slide.shapes:
            # 只处理非标题的文本框
            if shape.has_text_frame and not shape == new_slide.shapes.title:
                text_frame = shape.text_frame
                text_frame.clear()  # 清除原有内容
                # 将要点内容作为项目符号列表添加到文本框中
                for point in slide.content.bullet_points:
                    p = text_frame.add_paragraph()
                    p.text = point
                    p.level = 0  # 项目符号的级别
                    LOG.debug(f"添加列表项: {point}")
                break

        # 插入图片
        if slide.content.image_path:
            image_full_path = os.path.join(os.getcwd(), slide.content.image_path)  # 构建图片的绝对路径
            if os.path.exists(image_full_path):
                # 插入图片到占位符中
                for shape in new_slide.placeholders:
                    if shape.placeholder_format.type == 18:  # 18 表示图片占位符
                        shape.insert_picture(image_full_path)
                        LOG.debug(f"插入图片: {image_full_path}")
                        break
            else:
                LOG.warning(f"图片路径 '{image_full_path}' 不存在，跳过此图片。")

    # 保存生成的 PowerPoint 文件
    prs.save(output_path)
    LOG.info(f"演示文稿已保存到 '{output_path}'")
